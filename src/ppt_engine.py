import re
from pptx import Presentation
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from scraper import get_scraped_lyrics, find_song_url_by_name
import os

def generate_ppt(song_query, save_dir=None):
    """
    Takes a song name or URL from the UI, scrapes, 
    and builds the PowerPoint presentation.
    """
    
    # 1. Determine the URL based on input type
    if song_query.lower().startswith("http"):
        found_url = song_query
    else:
        found_url = find_song_url_by_name(song_query)

    # 2. Check if we have a valid link
    if not found_url:
        raise Exception(f"Could not find any song at URL or matching '{song_query}'.")

    # 3. Fetch live data from scraper.py
    song_title, english_slides, tamil_slides = get_scraped_lyrics(found_url)

    # 4. Presentation Setup
    prsnt = Presentation()
    prsnt.slide_width = Inches(13.333)
    prsnt.slide_height = Inches(7.5)
    blank_layout = prsnt.slide_layouts[6]

    # 5. Margins
    left_margin = Inches(0.3)
    box_width = Inches(12.533)
    box_height = Inches(3.1)
    top_box_y = Inches(0.5)
    bottom_box_y = Inches(3.9)

    # 6. Generate slides
    for eng_lines, tam_lines in zip(english_slides, tamil_slides):
        slide = prsnt.slides.add_slide(blank_layout)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0, 0, 0)

        # English Box
        english_box = slide.shapes.add_textbox(left_margin, top_box_y, box_width, box_height)
        tf_eng = english_box.text_frame
        tf_eng.word_wrap = True
        tf_eng.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        for i, line in enumerate(eng_lines):
            p = tf_eng.paragraphs[0] if i == 0 else tf_eng.add_paragraph()
            p.text = line.upper()
            p.alignment = PP_ALIGN.CENTER
            p.font.name = 'Google Sans'
            p.font.size = Pt(36)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)

        # Tamil Box
        tamil_box = slide.shapes.add_textbox(left_margin, bottom_box_y, box_width, box_height)
        tf_tam = tamil_box.text_frame
        tf_tam.word_wrap = True
        tf_tam.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        for i, line in enumerate(tam_lines):
            p = tf_tam.paragraphs[0] if i == 0 else tf_tam.add_paragraph()
            p.text = line
            p.alignment = PP_ALIGN.CENTER
            p.font.name = 'Arial Unicode MS'
            p.font.size = Pt(36)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)

    # 7. Save Logic
    safe_title = re.sub(r'[\\/*?:<>|]', "", song_title)
    filename = f"{safe_title}.pptx"

    target_folder = save_dir if (save_dir and os.path.isdir(save_dir)) else r"C:\Users\CHURCH OF I AM\Documents\pptautomation\Generated presentations"
    
    if not os.path.exists(target_folder):
        os.makedirs(target_folder)
        
    output_filepath = os.path.join(target_folder, filename)
    prsnt.save(output_filepath)
    
    return output_filepath